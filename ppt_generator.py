from pptx import Presentation

def generate_ppt_from_bullets(title, bullet_points, output_file="output\\generated.pptx"):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title

    bullets = bullet_points.split("\n")
    for i in range(0, len(bullets), 4):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Part {i//4 + 1}"
        content = slide.placeholders[1]
        content.text = "\n".join(f"• {point.strip()}" for point in bullets[i:i+4])

    prs.save(output_file)
    print(f"[✅] PPT saved to {output_file}")
