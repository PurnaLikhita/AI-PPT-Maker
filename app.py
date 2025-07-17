import streamlit as st
import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from datetime import datetime

# 🎨 Color Themes
THEMES = {
    "Classic Blue": {"bg": (255, 255, 255), "title": (0, 102, 204), "bullet": (0, 0, 0), "accent": (0, 51, 153)},
    "Dark Orange": {"bg": (30, 30, 30), "title": (255, 140, 0), "bullet": (200, 200, 200), "accent": (255, 85, 0)},
    "Elegant Purple": {"bg": (245, 240, 255), "title": (102, 0, 204), "bullet": (51, 0, 102), "accent": (204, 153, 255)},
    "Fresh Green": {"bg": (230, 255, 240), "title": (0, 153, 51), "bullet": (34, 85, 51), "accent": (0, 204, 102)},
    "Minimal Gray": {"bg": (250, 250, 250), "title": (60, 60, 60), "bullet": (90, 90, 90), "accent": (120, 120, 120)}
}

# 🔮 Mistral API
def call_mistral(prompt: str):
    response = requests.post(
        "http://localhost:11434/api/generate",
        json={"model": "mistral", "prompt": prompt, "stream": False}
    )
    resp_json = response.json()
    if "response" not in resp_json:
        raise Exception(f"API error: {resp_json}")
    return resp_json["response"]

# 🎨 Theme Styling
def apply_theme(slide, title_box, content_box, theme_colors):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*theme_colors["bg"])
    if title_box:
        para = title_box.text_frame.paragraphs[0]
        para.font.size = Pt(21)
        para.font.bold = True
        para.font.color.rgb = RGBColor(*theme_colors["title"])
    if content_box:
        for para in content_box.text_frame.paragraphs:
            para.font.size = Pt(20)
            para.font.color.rgb = RGBColor(*theme_colors["bullet"])

# 🧱 PPTX Creation
def create_presentation_from_text(text, filename="presentation.pptx", theme_colors=None):
    prs = Presentation()
    slides_text = text.strip().split("Slide")[1:]
    for slide_raw in slides_text:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        lines = slide_raw.strip().split("\n")
        if not lines:
            continue
        title = lines[0].strip(":").strip()
        bullets = [line.lstrip("-•").strip() for line in lines[1:] if line.strip()]
        title_box = slide.shapes.title
        content_box = slide.shapes.placeholders[1]
        font_size = Pt(18)
        if len(bullets) > 5:
            font_size = Pt(16)
        if len(bullets) > 8:
            font_size = Pt(14)
        tf = content_box.text_frame
        tf.clear()
        tf.word_wrap = True
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = font_size
        title_box.text = title
        if theme_colors:
            apply_theme(slide, title_box, content_box, theme_colors)
    prs.save(filename)
    return filename

# 🖼️ Streamlit UI
st.set_page_config(page_title="🧠 AI PPT Generator", layout="centered")
st.title("📊 AI PPT Generator with Mistral")
st.markdown("Offline slide generation using **Mistral via Ollama** 🧠💻")

# 🧠 Mode Selection
use_custom_prompt = st.checkbox("📝 Use Custom Prompt")

# 🎛️ Inputs
custom_prompt = ""
topic, description = "", ""
num_slides = 5

if use_custom_prompt:
    custom_prompt = st.text_area("✍️ Custom Prompt", height=200, value="Create a 5-slide presentation on Quantum Computing with slide titles and 4 bullet points each.")
else:
    topic = st.text_input("📌 Topic", "Applications of AI in Schools")
    description = st.text_area("📄 Description", "Example: Focus on hospital use cases like diagnosis, monitoring, etc.")
    num_slides = st.slider("📑 Number of Slides", 3, 10, 5)

theme_choice = st.selectbox("🎨 Theme", list(THEMES.keys()))
theme_colors = THEMES[theme_choice]

# 🚀 Generate
if st.button("🚀 Generate Presentation"):
    with st.spinner("Generating slides with Mistral..."):
        try:
            if use_custom_prompt:
                prompt = custom_prompt
            else:
                prompt = (
                    f"Create a {num_slides}-slide PowerPoint presentation on the topic: '{topic}'. "
                    f"Each slide should have a title and 3-5 bullet points. "
                    f"Structure it for clear presentation. Context: {description}"
                )
            text = call_mistral(prompt)
            filename = f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            create_presentation_from_text(text, filename, theme_colors)
            with open(filename, "rb") as f:
                st.success("✅ Presentation created!")
                st.download_button("📥 Download PPTX", f, file_name=filename, mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        except Exception as e:
            st.error(f"❌ Error: {e}")
